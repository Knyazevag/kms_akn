#!/usr/bin/env python3
"""
doc_to_obsidian.py — Интеграция RAG-системы с Obsidian.
Поддерживаемые форматы: PDF, DOCX, DOC, TXT, MD, XLSX, CSV, PPTX, ODT.

Для каждого документа из KMS/archive/ (рекурсивно):
  1. Извлекает текст первых страниц/блоков через соответствующий парсер.
  2. Отправляет текст в Ollama HTTP API для анализа.
  3. Создаёт .md файл в KMS/notes/ с YAML frontmatter.
  4. После обработки всех файлов — второй проход для построения wikilinks.

Запуск:
  python doc_to_obsidian.py [--doc-dir PATH] [--vault-dir PATH]
                             [--force] [--dry-run] [--model MODEL]
                             [--force-single PATH]
"""

from __future__ import annotations

import argparse
import copy
import json
import logging
import re
import sys
import unicodedata
from datetime import date
from pathlib import Path
from typing import Any

import requests
from tqdm import tqdm

# ─── Импорт конфигурации ───────────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import config
except ImportError as exc:
    sys.exit(f"[ERROR] Не удалось импортировать config.py: {exc}")

# ─── Константы ────────────────────────────────────────────────────────────
MAX_PAGES: int = 3            # сколько страниц/блоков читать для анализа
MAX_RETRIES: int = 3          # попыток запроса к Ollama (часть может уходить на рост бюджета)
MIN_COMMON_TAGS: int = config.WIKILINK_MIN_COMMON_TAGS
TEXT_LIMIT: int = 6000        # лимит символов текста, передаваемого в промпт
# Бюджет токенов на ПОЛНЫЙ JSON-объект метаданных. Дефолтный num_predict=1024
# (config.OLLAMA_OPTIONS) подходит для чата, но богатые статьи дают вывод длиннее
# 1024 токенов → JSON обрывается без закрывающей "}" → «не найден JSON-объект».
JSON_NUM_PREDICT: int = 2048
JSON_NUM_PREDICT_MAX: int = 8192  # потолок при удвоении бюджета на ретраях

# ─── Логирование ──────────────────────────────────────────────────────────

def _setup_logging() -> logging.Logger:
    logger = logging.getLogger("doc_to_obsidian")
    logger.setLevel(getattr(logging, config.LOG_LEVEL, logging.INFO))
    fmt = logging.Formatter(fmt=config.LOG_FORMAT, datefmt=config.LOG_DATE_FORMAT)
    fh = logging.FileHandler(config.LOG_FILE, encoding="utf-8")
    fh.setFormatter(fmt)
    logger.addHandler(fh)
    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(fmt)
    logger.addHandler(sh)
    return logger

log = _setup_logging()


# ══════════════════════════════════════════════════════════════════════════════
# Извлечение текста из разных форматов
# ══════════════════════════════════════════════════════════════════════════════

def _extract_preview_pdf(path: Path, max_pages: int = MAX_PAGES) -> str:
    """PDF: первые max_pages страниц. Автоматически запускает OCR для сканов."""
    try:
        import fitz
        import subprocess

        doc = fitz.open(str(path))
        # Быстрая проверка: есть ли текст на первых страницах?
        sample = sum(
            len(doc.load_page(i).get_text("text").strip())
            for i in range(min(3, len(doc)))
        )
        doc.close()

        # Если текста мало — возможно скан; пробуем OCR-версию
        if sample < 50:
            ocr_path = path.with_stem(path.stem + "_ocr")
            if ocr_path.exists():
                return _extract_preview_pdf(ocr_path, max_pages)
            # Проверяем доступность ocrmypdf
            try:
                chk = subprocess.run(
                    ["ocrmypdf", "--version"], capture_output=True, timeout=5
                )
                if chk.returncode == 0:
                    log.warning(
                        "PDF '%s' — скан без текста, запускаю OCR для превью...",
                        path.name
                    )
                    subprocess.run(
                        ["ocrmypdf", "--language", "rus+eng",
                         "--output-type", "pdf", "--skip-big",
                         "--optimize", "0", "--quiet",
                         str(path), str(ocr_path)],
                        capture_output=True, timeout=300
                    )
                    if ocr_path.exists():
                        return _extract_preview_pdf(ocr_path, max_pages)
            except (FileNotFoundError, subprocess.TimeoutExpired):
                pass
            log.warning(
                "PDF '%s' — скан без OCR-слоя, превью недоступно. "
                "Установите: sudo apt install ocrmypdf tesseract-ocr tesseract-ocr-rus",
                path.name
            )
            return ""

        # Обычный PDF с текстом
        doc = fitz.open(str(path))
        texts = []
        for i in range(min(max_pages, len(doc))):
            texts.append(doc.load_page(i).get_text("text"))
        doc.close()
        return "\n".join(texts)
    except ImportError:
        log.error("PyMuPDF не установлен: pip install pymupdf")
        return ""
    except Exception as e:
        log.warning("Ошибка чтения PDF %s: %s", path.name, e)
        return ""


def _extract_preview_docx(path: Path, max_pages: int = MAX_PAGES) -> str:
    """DOCX: первые ~max_pages*500 символов из параграфов."""
    try:
        from docx import Document
        doc = Document(str(path))
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(lines[:max_pages * 20])   # ~20 абзацев = 1 «страница»
    except ImportError:
        log.error("python-docx не установлен: pip install python-docx")
        return ""
    except Exception as e:
        log.warning("Ошибка чтения DOCX %s: %s", path.name, e)
        return ""


def _extract_preview_txt(path: Path, max_pages: int = MAX_PAGES) -> str:
    """TXT/MD: первые TEXT_LIMIT символов."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text[:TEXT_LIMIT * max_pages]
    except Exception as e:
        log.warning("Ошибка чтения TXT %s: %s", path.name, e)
        return ""


def _extract_preview_xlsx(path: Path, max_pages: int = MAX_PAGES) -> str:
    """XLSX: первые max_pages листов, первые 30 строк каждого."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        result = []
        for sheet_name in list(wb.sheetnames)[:max_pages]:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 30:
                    break
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                result.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(result)
    except ImportError:
        log.error("openpyxl не установлен: pip install openpyxl")
        return ""
    except Exception as e:
        log.warning("Ошибка чтения XLSX %s: %s", path.name, e)
        return ""


def _extract_preview_csv(path: Path, max_pages: int = MAX_PAGES) -> str:
    """CSV: первые 50 строк."""
    try:
        import csv
        rows = []
        for enc in ("utf-8", "cp1251", "latin-1"):
            try:
                with open(path, newline="", encoding=enc, errors="replace") as f:
                    reader = csv.reader(f)
                    for i, row in enumerate(reader):
                        if i >= 50:
                            break
                        rows.append(" | ".join(row))
                break
            except UnicodeDecodeError:
                continue
        return "\n".join(rows)
    except Exception as e:
        log.warning("Ошибка чтения CSV %s: %s", path.name, e)
        return ""


def _extract_preview_pptx(path: Path, max_pages: int = MAX_PAGES) -> str:
    """PPTX: текст первых max_pages слайдов."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        result = []
        for i, slide in enumerate(prs.slides):
            if i >= max_pages * 3:   # берём больше слайдов для контекста
                break
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                result.append(f"[Слайд {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(result)
    except ImportError:
        log.error("python-pptx не установлен: pip install python-pptx")
        return ""
    except Exception as e:
        log.warning("Ошибка чтения PPTX %s: %s", path.name, e)
        return ""


def _extract_preview_odt(path: Path, max_pages: int = MAX_PAGES) -> str:
    """ODT: первые абзацы документа."""
    try:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype
        doc = load(str(path))
        paragraphs = doc.text.getElementsByType(P)
        lines = [teletype.extractText(p) for p in paragraphs if teletype.extractText(p).strip()]
        return "\n".join(lines[:max_pages * 20])
    except ImportError:
        log.error("odfpy не установлен: pip install odfpy")
        return ""
    except Exception as e:
        log.warning("Ошибка чтения ODT %s: %s", path.name, e)
        return ""


# Диспетчер: расширение → функция извлечения превью
PREVIEW_EXTRACTORS = {
    ".pdf":  _extract_preview_pdf,
    ".docx": _extract_preview_docx,
    ".doc":  _extract_preview_docx,   # python-docx читает и .doc через OOXML
    ".txt":  _extract_preview_txt,
    ".md":   _extract_preview_txt,
    ".xlsx": _extract_preview_xlsx,
    ".csv":  _extract_preview_csv,
    ".pptx": _extract_preview_pptx,
    ".odt":  _extract_preview_odt,
}


def extract_preview(path: Path) -> str:
    """Извлекает текст-превью документа для анализа через Ollama."""
    ext = path.suffix.lower()
    extractor = PREVIEW_EXTRACTORS.get(ext)
    if extractor is None:
        log.warning("Формат '%s' не поддерживается для превью: %s", ext, path.name)
        return ""
    return extractor(path)


# ══════════════════════════════════════════════════════════════════════════════
# Вспомогательные функции
# ══════════════════════════════════════════════════════════════════════════════

def slugify(text: str) -> str:
    """'CO2 Storage: a Review' → 'co2_storage_a_review'"""
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    text = text.lower()
    text = re.sub(r"[^\w\s-]", " ", text)
    text = re.sub(r"[\s\-]+", "_", text).strip("_")
    return text or "untitled"


def normalise_tag(tag: str) -> str:
    tag = tag.strip().lower()
    return re.sub(r"[\s\-]+", "_", tag)


def _all_tags(meta: dict[str, Any]) -> list[str]:
    return list(meta["tags_from_taxonomy"]) + list(meta["tags_auto"])


def _yaml_list(items: list[str]) -> str:
    escaped = [f'"{i}"' for i in items]
    return "[" + ", ".join(escaped) + "]"


# ══════════════════════════════════════════════════════════════════════════════
# Взаимодействие с Ollama
# ══════════════════════════════════════════════════════════════════════════════

def _build_prompt(text: str, taxonomy: list[str], file_type: str) -> str:
    taxonomy_str = ", ".join(taxonomy)
    truncated = text[:TEXT_LIMIT]
    type_hint = {
        "pdf": "scientific PDF article",
        "docx": "Word document",
        "doc": "Word document",
        "txt": "plain text document",
        "md": "Markdown document",
        "xlsx": "Excel spreadsheet",
        "csv": "CSV data file",
        "pptx": "PowerPoint presentation",
        "odt": "LibreOffice text document",
    }.get(file_type, "document")

    return f"""You are a scientific metadata extraction assistant for petroleum engineering documents.
Analyze the following text extracted from a {type_hint}.

TAXONOMY OF ALLOWED TAGS (choose ONLY from this list for tags_from_taxonomy):
{taxonomy_str}

If you believe an important concept is NOT covered by the taxonomy, add it to tags_auto
with the prefix "auto/" (e.g. "auto/supercritical_co2"). Keep auto tags minimal.

All tags must be lowercase with spaces replaced by underscores.

Return ONLY a valid JSON object — no markdown, no explanations, no code fences.
Use null (not "null") for missing fields. Use empty arrays [] for missing lists.

Required JSON structure:
{{
  "title": "Document title in original language",
  "authors": ["Author1", "Author2"],
  "year": 2023,
  "journal": "Journal, conference, or source name (null if unknown)",
  "doi": "10.xxxx/xxxx or null",
  "keywords": ["keyword1", "keyword2"],
  "tags_from_taxonomy": ["tag1", "tag2", "tag3"],
  "tags_auto": ["auto/new_concept"],
  "summary_ru": "Краткое резюме на русском языке (3-5 предложений).",
  "summary_en": "Brief summary in English (2-3 sentences).",
  "methods": ["method1", "method2"],
  "software": ["PHREEQC", "tNavigator"],
  "key_findings": ["finding1", "finding2", "finding3"]
}}

DOCUMENT TEXT:
{truncated}
"""


def call_ollama(
    text: str,
    model: str,
    taxonomy: list[str],
    file_type: str = "pdf",
    max_retries: int = MAX_RETRIES,
) -> dict[str, Any] | None:
    prompt = _build_prompt(text, taxonomy, file_type)
    # Локальная копия опций: не трогаем глобальный config.OLLAMA_OPTIONS (его
    # использует и чат). Берём бюджет не меньше JSON_NUM_PREDICT.
    options = copy.deepcopy(config.OLLAMA_OPTIONS)
    num_predict = max(int(options.get("num_predict", 0)), JSON_NUM_PREDICT)

    for attempt in range(1, max_retries + 1):
        options["num_predict"] = num_predict
        payload = {
            "model": model,
            "prompt": prompt,
            "stream": False,
            "think": False,   # thinking-модели (qwen3/r1) иначе не вернут JSON в response
            "format": "json",  # форсируем JSON-грамматику: без прозы и ```-ограждений
            "options": options,
        }
        try:
            response = requests.post(
                config.OLLAMA_API_URL,
                json=payload,
                timeout=config.OLLAMA_TIMEOUT,
            )
            response.raise_for_status()
            data = response.json()
            raw = data.get("response", "")
            # Обрыв по лимиту токенов → JSON неполный. Удваиваем бюджет и
            # повторяем (ретрай с теми же параметрами был бы бесполезен).
            if data.get("done_reason") == "length" and num_predict < JSON_NUM_PREDICT_MAX:
                log.warning(
                    "Попытка %d/%d: вывод обрезан по num_predict=%d — увеличиваю бюджет",
                    attempt, max_retries, num_predict,
                )
                num_predict = min(num_predict * 2, JSON_NUM_PREDICT_MAX)
                continue
            json_match = re.search(r"\{[\s\S]*\}", raw)
            if not json_match:
                raise ValueError("В ответе Ollama не найден JSON-объект")
            return json.loads(json_match.group())
        except (requests.RequestException, ValueError, json.JSONDecodeError) as exc:
            log.warning("Попытка %d/%d: ошибка Ollama — %s", attempt, max_retries, exc)

    return None


# ══════════════════════════════════════════════════════════════════════════════
# Формирование содержимого заметки
# ══════════════════════════════════════════════════════════════════════════════

EMPTY_META: dict[str, Any] = {
    "title": "", "authors": [], "year": None, "journal": "",
    "doi": None, "keywords": [], "tags_from_taxonomy": [], "tags_auto": [],
    "summary_ru": "", "summary_en": "", "methods": [], "software": [],
    "key_findings": [],
}


def _normalise_metadata(raw: dict[str, Any], file_name: str) -> dict[str, Any]:
    meta = {**EMPTY_META, **raw}

    if not meta.get("title"):
        meta["title"] = Path(file_name).stem

    tax_tags = [normalise_tag(t) for t in (meta.get("tags_from_taxonomy") or [])]
    auto_tags_raw = meta.get("tags_auto") or []
    auto_tags: list[str] = []
    for t in auto_tags_raw:
        t = normalise_tag(t)
        if not t.startswith("auto/"):
            t = "auto/" + t
        auto_tags.append(t)

    meta["tags_from_taxonomy"] = tax_tags
    meta["tags_auto"] = auto_tags

    for key in ("authors", "keywords", "methods", "software", "key_findings"):
        if not isinstance(meta.get(key), list):
            meta[key] = [str(meta[key])] if meta.get(key) else []

    return meta


def _build_note_content(
    meta: dict[str, Any],
    source_filename: str,
    file_type: str,
    today: str,
) -> str:
    title = meta["title"]
    authors = meta["authors"]
    year = meta.get("year") or "н/д"
    journal = meta.get("journal") or "н/д"
    doi = meta.get("doi") or "н/д"
    summary_ru = meta.get("summary_ru") or ""
    summary_en = meta.get("summary_en") or ""
    methods = meta.get("methods") or []
    software = meta.get("software") or []
    key_findings = meta.get("key_findings") or []
    tags = _all_tags(meta)

    # Тип файла в человекочитаемом виде
    type_labels = {
        "pdf": "PDF", "docx": "Word (DOCX)", "doc": "Word (DOC)",
        "txt": "Текстовый файл", "md": "Markdown",
        "xlsx": "Excel (XLSX)", "csv": "CSV",
        "pptx": "PowerPoint (PPTX)", "odt": "LibreOffice (ODT)",
    }
    type_label = type_labels.get(file_type, file_type.upper())

    # YAML frontmatter
    frontmatter = f"""---
title: "{title}"
authors: {_yaml_list(authors)}
year: {year}
journal: "{journal}"
doi: "{doi}"
tags: {_yaml_list(tags)}
date_indexed: {today}
source_file: "{source_filename}"
file_type: "{file_type}"
---"""

    authors_str = ", ".join(authors) if authors else "н/д"
    methods_str = ", ".join(methods) if methods else "н/д"
    software_str = ", ".join(software) if software else "н/д"
    findings_md = "\n".join(f"- {f}" for f in key_findings) if key_findings else "- н/д"
    tags_obsidian = " ".join(f"#{t}" for t in tags) if tags else ""

    body = f"""
# {title}

## Метаданные
| Поле | Значение |
|------|----------|
| Тип документа | {type_label} |
| Авторы | {authors_str} |
| Год | {year} |
| Журнал/Источник | {journal} |
| DOI | {doi} |

## Резюме
{summary_ru}

*English summary: {summary_en}*

## Методы и ПО
**Методы:** {methods_str}
**Программное обеспечение:** {software_str}

## Ключевые выводы
{findings_md}

## Теги
{tags_obsidian}

## Связанные документы
*Будет заполнено автоматически при наличии общих тегов*

---
*Источник: {source_filename} | Проиндексировано: {today}*
"""
    return frontmatter + "\n" + body


# ══════════════════════════════════════════════════════════════════════════════
# Второй проход: wikilinks
# ══════════════════════════════════════════════════════════════════════════════

RELATED_SECTION_MARKER = "## Связанные документы"
RELATED_PLACEHOLDER = "*Будет заполнено автоматически при наличии общих тегов*"


def _find_related(
    note_data: list[dict[str, Any]],
    min_common: int = MIN_COMMON_TAGS,
) -> dict[str, list[tuple[str, set[str]]]]:
    relations: dict[str, list[tuple[str, set[str]]]] = {
        nd["slug"]: [] for nd in note_data
    }
    for i, a in enumerate(note_data):
        for b in note_data[i + 1:]:
            common = set(a["tags"]) & set(b["tags"])
            if len(common) >= min_common:
                relations[a["slug"]].append((b["title"], common))
                relations[b["slug"]].append((a["title"], common))
    return relations


def _update_related_section(
    md_path: Path,
    related: list[tuple[str, set[str]]],
    dry_run: bool = False,
) -> None:
    if not related:
        return

    content = md_path.read_text(encoding="utf-8")
    lines = []
    for title, common_tags in sorted(related, key=lambda x: -len(x[1])):
        tags_str = " ".join(f"#{t}" for t in sorted(common_tags))
        lines.append(f"- [[{title}]] ({len(common_tags)} общих тег(а): {tags_str})")
    links_block = "\n".join(lines)

    new_section = (
        f"{RELATED_SECTION_MARKER}\n"
        f"{RELATED_PLACEHOLDER}\n"
        f"{links_block}"
    )

    if RELATED_PLACEHOLDER in content:
        new_content = content.replace(
            f"{RELATED_SECTION_MARKER}\n{RELATED_PLACEHOLDER}",
            new_section,
        )
    else:
        pattern = re.compile(
            rf"({re.escape(RELATED_SECTION_MARKER)}\n).*?(\n---)",
            re.DOTALL,
        )
        new_content = pattern.sub(
            lambda m: m.group(1) + links_block + "\n" + m.group(2),
            content,
        )

    if dry_run:
        log.info("[DRY-RUN] Обновление связей в: %s", md_path.name)
        return

    md_path.write_text(new_content, encoding="utf-8")


# ══════════════════════════════════════════════════════════════════════════════
# Основная логика
# ══════════════════════════════════════════════════════════════════════════════

def process_archive(
    doc_dir: Path,
    vault_dir: Path,
    model: str,
    force: bool = False,
    dry_run: bool = False,
    single_file: Path | None = None,
) -> None:
    today = date.today().isoformat()
    taxonomy: list[str] = config.TAXONOMY

    # Определяем список файлов
    if single_file:
        # Режим --force-single: обрабатываем один файл
        if single_file.suffix.lower() not in config.SUPPORTED_EXTENSIONS:
            log.warning("Формат '%s' не поддерживается: %s",
                        single_file.suffix, single_file.name)
            return
        files = [single_file]
        log.info("Режим одиночного файла: %s", single_file.name)
    else:
        # Полный обход директории — все поддерживаемые форматы рекурсивно
        files = []
        for ext in config.SUPPORTED_EXTENSIONS:
            files.extend(doc_dir.rglob(f"*{ext}"))
        files = sorted(set(files))

    if not files:
        log.warning("Файлы не найдены в %s", doc_dir)
        return

    log.info("Найдено файлов для обработки: %d", len(files))

    if not dry_run:
        vault_dir.mkdir(parents=True, exist_ok=True)

    # Уже обработанные файлы (по полю source_file во фронтматтере заметок).
    # Имя заметки зависит от заголовка, который выдаёт LLM, поэтому проверки
    # по имени файла-заметки недостаточно — без этого повторные прогоны
    # (или авто-обработка) плодили бы дубли. --force отключает пропуск.
    existing_sources: set[str] = set()
    if not force:
        for _md in vault_dir.glob("*.md"):
            _head = _md.read_text(encoding="utf-8", errors="ignore")[:2000]
            _m = re.search(r'^source_file:\s*"?(.+?)"?\s*$', _head, re.MULTILINE)
            if _m:
                existing_sources.add(_m.group(1).strip())

    note_data: list[dict[str, Any]] = []

    # ── Первый проход: создание заметок ───────────────────────────────────
    for file_path in tqdm(files, desc="Создание заметок", unit="файл"):
        file_name = file_path.name
        file_type = file_path.suffix.lower().lstrip(".")

        if (not force) and file_name in existing_sources:
            log.info("Заметка для %s уже существует — пропуск", file_name)
            continue

        log.info("Обрабатываю [%s]: %s", file_type.upper(), file_name)

        # Извлечение превью
        text = extract_preview(file_path)
        if not text.strip():
            log.warning("Текст не извлечён из %s — пропуск", file_name)
            continue

        # Анализ через Ollama
        if dry_run:
            log.info("[DRY-RUN] Пропуск вызова Ollama для %s", file_name)
            raw_meta: dict[str, Any] = {}
            failed = False
        else:
            raw_meta = call_ollama(text, model=model, taxonomy=taxonomy, file_type=file_type)
            failed = raw_meta is None
            if failed:
                log.error("Ollama не вернула валидный JSON для %s", file_name)
                raw_meta = {}

        meta = _normalise_metadata(raw_meta, file_name)
        slug = slugify(meta["title"])
        md_path = vault_dir / (slug + ".md")

        # Пропускаем существующие (если не --force)
        if md_path.exists() and not force:
            log.info("Заметка существует, пропуск: %s", md_path.name)
            existing = md_path.read_text(encoding="utf-8")
            existing_tags = re.findall(r'"([\w/]+)"', existing[:1000])
            note_data.append({
                "slug": slug, "title": meta["title"],
                "tags": existing_tags, "md_path": md_path,
            })
            continue

        note_content = _build_note_content(meta, file_name, file_type, today)

        if failed:
            note_content = note_content.replace(
                f"# {meta['title']}",
                f"# [!] {meta['title']}\n\n"
                "> **[!] Требует ручной проверки** — Ollama не вернула метаданные.",
            )

        if dry_run:
            log.info("[DRY-RUN] Создание заметки: %s", md_path.name)
            print(f"\n{'─'*60}\n[DRY-RUN] {md_path.name}\n{'─'*60}")
            print(note_content[:500] + "...\n")
        else:
            md_path.write_text(note_content, encoding="utf-8")
            log.info("Заметка создана: %s", md_path.name)

        note_data.append({
            "slug": slug, "title": meta["title"],
            "tags": _all_tags(meta), "md_path": md_path,
        })

    # ── Второй проход: wikilinks ──────────────────────────────────────────
    log.info("Второй проход: построение wikilinks...")
    relations = _find_related(note_data, min_common=MIN_COMMON_TAGS)
    updated = 0
    for nd in tqdm(note_data, desc="Wikilinks", unit="заметка"):
        related = relations.get(nd["slug"], [])
        if not related:
            continue
        md_path = nd["md_path"]
        if not md_path.exists() and not dry_run:
            continue
        _update_related_section(md_path, related, dry_run=dry_run)
        updated += 1

    log.info(
        "Готово. Обработано файлов: %d. Заметок с wikilinks: %d.",
        len(note_data), updated,
    )


# ══════════════════════════════════════════════════════════════════════════════
# CLI
# ══════════════════════════════════════════════════════════════════════════════

def parse_args() -> argparse.Namespace:
    exts = ", ".join(sorted(config.SUPPORTED_EXTENSIONS))
    parser = argparse.ArgumentParser(
        description=f"Генерация Obsidian-заметок из документов ({exts}) через Ollama.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--doc-dir", "--pdf-dir",   # --pdf-dir для обратной совместимости
        type=Path,
        default=config.KMS_ARCHIVE_DIR,
        help="Путь к архиву документов (default: %(default)s)",
    )
    parser.add_argument(
        "--vault-dir",
        type=Path,
        default=config.KMS_VAULT_DIR,
        help="Путь к Obsidian vault (default: %(default)s)",
    )
    parser.add_argument("--force", action="store_true",
                        help="Перезаписать существующие заметки")
    parser.add_argument("--dry-run", action="store_true",
                        help="Показать что будет сделано без записи файлов")
    parser.add_argument("--model", type=str, default=None,
                        help="Модель Ollama (override config.OLLAMA_MODEL)")
    parser.add_argument("--force-single", type=Path, default=None,
                        metavar="FILE",
                        help="Обработать один конкретный файл (используется watcher)")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    model = args.model or config.OLLAMA_MODEL

    log.info("=" * 60)
    log.info("doc_to_obsidian.py — старт")
    log.info("Архив       : %s", args.doc_dir)
    log.info("Vault       : %s", args.vault_dir)
    log.info("Модель      : %s", model)
    log.info("--force     : %s", args.force)
    log.info("--dry-run   : %s", args.dry_run)
    log.info("--force-single: %s", args.force_single)
    log.info("=" * 60)

    doc_dir = Path(args.doc_dir)
    vault_dir = Path(args.vault_dir)

    if not doc_dir.exists() and not args.force_single:
        log.error("Директория архива не найдена: %s", doc_dir)
        sys.exit(1)

    process_archive(
        doc_dir=doc_dir,
        vault_dir=vault_dir,
        model=model,
        force=args.force,
        dry_run=args.dry_run,
        single_file=args.force_single,
    )


if __name__ == "__main__":
    main()
