"""
ingest.py — Индексация документов в векторную базу ChromaDB.

Поддерживаемые форматы: PDF, DOCX, DOC, TXT, MD, XLSX, CSV, PPTX, ODT.
Список форматов управляется через config.SUPPORTED_EXTENSIONS.

Алгоритм:
  1. Рекурсивный обход директории архива.
  2. Определение формата файла и выбор парсера.
  3. Нарезка текста на чанки с перекрытием.
  4. Батчевая генерация эмбеддингов через sentence-transformers.
  5. Сохранение чанков + эмбеддингов + метаданных в ChromaDB.

Запуск:
  python ingest.py [--doc-dir /path/to/docs] [--reset]
"""

import argparse
import hashlib
import json
import logging
import re
import subprocess
import sys
import time
from pathlib import Path
from typing import Generator

import fitz  # PyMuPDF — для PDF
from langdetect import detect, LangDetectException
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
import chromadb
from chromadb.config import Settings

import config

# ─────────────────────────────────────────────────────────────
# Настройка логирования
# ─────────────────────────────────────────────────────────────

logging.basicConfig(
    level=getattr(logging, config.LOG_LEVEL),
    format=config.LOG_FORMAT,
    datefmt=config.LOG_DATE_FORMAT,
    handlers=[
        logging.FileHandler(config.LOG_FILE, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)
logger = logging.getLogger("ingest")


# ─────────────────────────────────────────────────────────────
# Парсеры для каждого формата
# ─────────────────────────────────────────────────────────────

def _is_ocrmypdf_available() -> bool:
    """Проверяет, установлен ли ocrmypdf."""
    try:
        result = subprocess.run(
            ["ocrmypdf", "--version"],
            capture_output=True, text=True, timeout=5
        )
        return result.returncode == 0
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False


def _run_ocr_on_pdf(path: Path) -> "Path | None":
    """
    Запускает ocrmypdf на скане без текста.
    Создаёт файл <имя>_ocr.pdf рядом с оригиналом.
    Возвращает путь к OCR-версии или None при ошибке.
    """
    ocr_path = path.with_stem(path.stem + "_ocr")
    if ocr_path.exists():
        logger.info(f"OCR: найден кешированный файл '{ocr_path.name}', использую его.")
        return ocr_path

    logger.warning(
        f"PDF '{path.name}' не содержит текстового слоя — скан без OCR. "
        f"Запускаю ocrmypdf (rus+eng)..."
    )
    try:
        result = subprocess.run(
            [
                "ocrmypdf",
                "--language", "rus+eng",
                "--output-type", "pdf",
                "--skip-big",      # пропускать изображения > 4 МПикс (защита от зависания)
                "--optimize", "0", # не оптимизировать — ускоряет работу
                "--quiet",
                str(path),
                str(ocr_path),
            ],
            capture_output=True, text=True, timeout=300  # макс. 5 минут
        )
        if result.returncode == 0 and ocr_path.exists():
            logger.info(f"OCR: успешно создан '{ocr_path.name}'")
            return ocr_path
        else:
            logger.error(
                f"ocrmypdf завершился с ошибкой (код {result.returncode}) "
                f"для '{path.name}': {result.stderr[:300]}"
            )
            return None
    except subprocess.TimeoutExpired:
        logger.error(f"OCR: превышено время ожидания (5 мин) для '{path.name}'")
        return None
    except FileNotFoundError:
        logger.error(
            "ocrmypdf не установлен. Установите: "
            "sudo apt install ocrmypdf tesseract-ocr tesseract-ocr-rus"
        )
        return None


def _pdf_has_text(path: Path, sample_pages: int = 3) -> bool:
    """
    Быстрая проверка: содержит ли PDF текстовый слой.
    Проверяет первые sample_pages страниц.
    Возвращает True, если суммарно найдено >= 50 символов текста.
    """
    try:
        doc = fitz.open(str(path))
        total_chars = 0
        pages_to_check = min(sample_pages, len(doc))
        for i in range(pages_to_check):
            total_chars += len(doc[i].get_text("text").strip())
        doc.close()
        return total_chars >= 50
    except Exception:
        return False


def extract_pdf(path: Path) -> list[dict]:
    """PDF → список {'page_num': int, 'text': str}.

    Автоматически обнаруживает сканы без OCR-слоя и запускает ocrmypdf.
    Если ocrmypdf не установлен — логирует предупреждение и возвращает [].
    """
    pages = []

    # Проверяем наличие текстового слоя
    if not _pdf_has_text(path):
        if _is_ocrmypdf_available():
            ocr_path = _run_ocr_on_pdf(path)
            if ocr_path is not None:
                # Рекурсивно извлекаем текст из OCR-версии
                pages = extract_pdf(ocr_path)
                if pages:
                    logger.info(
                        f"OCR: '{path.name}' — извлечено {len(pages)} страниц "
                        f"после распознавания."
                    )
                return pages
        else:
            logger.warning(
                f"PDF '{path.name}' является сканом без текстового слоя. "
                f"Файл пропущен. Для автоматического OCR установите: "
                f"sudo apt install ocrmypdf tesseract-ocr tesseract-ocr-rus"
            )
            return []

    # Обычный PDF с текстом — стандартное извлечение
    try:
        doc = fitz.open(str(path))
        for page_num, page in enumerate(doc, start=1):
            try:
                text = _clean_text(page.get_text("text"))
                if text:
                    pages.append({"page_num": page_num, "text": text})
            except Exception as e:
                logger.warning(f"Ошибка страницы {page_num} в {path.name}: {e}")
        doc.close()
        logger.info(f"PDF: извлечено {len(pages)} страниц из '{path.name}'")
    except Exception as e:
        logger.error(f"Не удалось открыть PDF '{path}': {e}")
    return pages


def extract_docx(path: Path) -> list[dict]:
    """DOCX → список страниц (разбивка по параграфам, ~500 символов)"""
    try:
        from docx import Document  # python-docx
        doc = Document(str(path))
        full_text = "\n".join(
            p.text for p in doc.paragraphs if p.text.strip()
        )
        # Добавляем текст из таблиц
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text += "\n" + row_text
        return _text_to_pages(full_text, path.name)
    except ImportError:
        logger.error("python-docx не установлен. Выполните: pip install python-docx")
        return []
    except Exception as e:
        logger.error(f"Ошибка чтения DOCX '{path}': {e}")
        return []


def extract_doc(path: Path) -> list[dict]:
    """DOC (старый формат Word) → список страниц через LibreOffice или antiword."""
    # Сначала пробуем конвертацию через LibreOffice (надёжнее)
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text",
             "--outdir", str(path.parent), str(path)],
            capture_output=True, text=True, timeout=30
        )
        txt_path = path.with_suffix(".txt")
        if txt_path.exists():
            pages = extract_txt(txt_path)
            txt_path.unlink()  # удаляем временный файл
            logger.info(f"DOC: конвертирован через LibreOffice '{path.name}'")
            return pages
    except FileNotFoundError:
        pass  # LibreOffice не установлен, пробуем antiword
    except Exception as e:
        logger.warning(f"LibreOffice не смог конвертировать '{path.name}': {e}")

    # Запасной вариант — antiword
    try:
        result = subprocess.run(
            ["antiword", str(path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            logger.info(f"DOC: прочитан через antiword '{path.name}'")
            return _text_to_pages(result.stdout, path.name)
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning(f"antiword не смог прочитать '{path.name}': {e}")

    logger.error(
        f"Не удалось прочитать DOC '{path.name}'. "
        "Установите LibreOffice: sudo apt install libreoffice"
    )
    return []


def extract_txt(path: Path) -> list[dict]:
    """TXT или MD → список страниц."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        logger.info(f"TXT/MD: прочитан '{path.name}' ({len(text)} символов)")
        return _text_to_pages(text, path.name)
    except Exception as e:
        logger.error(f"Ошибка чтения TXT '{path}': {e}")
        return []


def extract_xlsx(path: Path) -> list[dict]:
    """XLSX → список страниц (каждый лист — отдельный «блок текста»)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        all_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                all_text.append(f"[Лист: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        full_text = "\n\n".join(all_text)
        logger.info(f"XLSX: прочитано {len(wb.sheetnames)} листов из '{path.name}'")
        return _text_to_pages(full_text, path.name)
    except ImportError:
        logger.error("openpyxl не установлен. Выполните: pip install openpyxl")
        return []
    except Exception as e:
        logger.error(f"Ошибка чтения XLSX '{path}': {e}")
        return []


def extract_csv(path: Path) -> list[dict]:
    """CSV → список страниц."""
    try:
        import csv
        rows = []
        # Пробуем разные кодировки
        for encoding in ("utf-8", "cp1251", "latin-1"):
            try:
                with open(path, newline="", encoding=encoding, errors="replace") as f:
                    reader = csv.reader(f)
                    rows = [" | ".join(r) for r in reader if any(c.strip() for c in r)]
                break
            except UnicodeDecodeError:
                continue
        full_text = "\n".join(rows)
        logger.info(f"CSV: прочитано {len(rows)} строк из '{path.name}'")
        return _text_to_pages(full_text, path.name)
    except Exception as e:
        logger.error(f"Ошибка чтения CSV '{path}': {e}")
        return []


def extract_pptx(path: Path) -> list[dict]:
    """PPTX → список страниц (каждый слайд — отдельная «страница»)."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        pages = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                pages.append({
                    "page_num": slide_num,
                    "text": _clean_text("\n".join(texts))
                })
        logger.info(f"PPTX: извлечено {len(pages)} слайдов из '{path.name}'")
        return pages
    except ImportError:
        logger.error("python-pptx не установлен. Выполните: pip install python-pptx")
        return []
    except Exception as e:
        logger.error(f"Ошибка чтения PPTX '{path}': {e}")
        return []


def extract_odt(path: Path) -> list[dict]:
    """ODT (LibreOffice Writer) → список страниц."""
    try:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype
        doc = load(str(path))
        paragraphs = doc.text.getElementsByType(P)
        lines = [teletype.extractText(p) for p in paragraphs]
        full_text = "\n".join(line for line in lines if line.strip())
        logger.info(f"ODT: прочитан '{path.name}' ({len(lines)} параграфов)")
        return _text_to_pages(full_text, path.name)
    except ImportError:
        logger.error("odfpy не установлен. Выполните: pip install odfpy")
        return []
    except Exception as e:
        logger.error(f"Ошибка чтения ODT '{path}': {e}")
        return []


# ─────────────────────────────────────────────────────────────
# Диспетчер форматов
# ─────────────────────────────────────────────────────────────

EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".doc":  extract_doc,
    ".txt":  extract_txt,
    ".md":   extract_txt,
    ".xlsx": extract_xlsx,
    ".csv":  extract_csv,
    ".pptx": extract_pptx,
    ".odt":  extract_odt,
}


def extract_document(path: Path) -> list[dict]:
    """Универсальный экстрактор: выбирает парсер по расширению файла."""
    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        logger.warning(f"Формат '{ext}' не поддерживается: {path.name}")
        return []
    return extractor(path)


# ─────────────────────────────────────────────────────────────
# Вспомогательные функции
# ─────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """Нормализует пробелы и переносы строк."""
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _text_to_pages(text: str, filename: str, page_size: int = 2000) -> list[dict]:
    """
    Разбивает длинный текст на «страницы» по ~page_size символов.
    Используется для форматов без нативной нумерации страниц (TXT, DOCX и т.д.).
    Разбивка происходит по границам абзацев, а не посередине слова.
    """
    text = _clean_text(text)
    if not text:
        return []

    pages = []
    page_num = 1
    start = 0

    while start < len(text):
        end = start + page_size
        if end >= len(text):
            chunk = text[start:]
        else:
            # Ищем ближайший перенос строки, чтобы не резать посередине абзаца
            newline_pos = text.rfind("\n", start, end)
            if newline_pos > start + page_size // 2:
                end = newline_pos
            chunk = text[start:end]

        chunk = chunk.strip()
        if chunk:
            pages.append({"page_num": page_num, "text": chunk})
            page_num += 1

        start = end

    logger.debug(f"Текст из '{filename}' разбит на {len(pages)} страниц")
    return pages


def make_doc_id(file_name: str, page_num: int, chunk_idx: int) -> str:
    """Создаёт уникальный ID для чанка на основе имени файла и позиции."""
    raw = f"{file_name}::{page_num}::{chunk_idx}"
    return hashlib.md5(raw.encode()).hexdigest()


def detect_language(text: str) -> str:
    """Определяет язык текста ('ru', 'en', 'unknown')."""
    try:
        lang = detect(text[:500])
        return lang if lang in ("ru", "en") else "other"
    except LangDetectException:
        return "unknown"


def chunk_text(text: str) -> list[str]:
    """
    Нарезает текст на перекрывающиеся чанки.
    Параметры берутся из config.
    """
    chunks = []
    start = 0
    while start < len(text):
        end = start + config.CHUNK_SIZE
        chunk = text[start:end].strip()
        if len(chunk) >= config.MIN_CHUNK_LENGTH:
            chunks.append(chunk)
        start += config.CHUNK_SIZE - config.CHUNK_OVERLAP
    return chunks


# ─────────────────────────────────────────────────────────────
# ChromaDB — инициализация
# ─────────────────────────────────────────────────────────────

def get_collection(reset: bool = False):
    """Инициализирует ChromaDB и возвращает коллекцию."""
    client = chromadb.PersistentClient(
        path=str(config.CHROMA_PERSIST_DIR),
        settings=Settings(anonymized_telemetry=False),
    )

    if reset:
        try:
            client.delete_collection(config.CHROMA_COLLECTION_NAME)
            logger.info("Коллекция ChromaDB сброшена.")
        except Exception:
            pass

    collection = client.get_or_create_collection(
        name=config.CHROMA_COLLECTION_NAME,
        metadata={
            "hnsw:space": "cosine",
            "description": "Научные документы (PDF, DOCX, TXT и др.)",
        },
    )
    return collection


# ─────────────────────────────────────────────────────────────
# Основной генератор чанков
# ─────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────
# Манифест индексации: file-level skip неизменённых файлов
# ─────────────────────────────────────────────────────────────
# Хранит {путь_файла: "размер:mtime"} для файлов, чьи чанки уже в базе.
# Позволяет НЕ парсить (дорогая операция для больших PDF) файлы, которые не
# менялись с прошлой индексации. Манифест записывается один раз в конце успешного
# прогона; при сбое он не обновляется → файлы переиндексируются (безопасно).
MANIFEST_PATH = config.CHROMA_PERSIST_DIR / "ingest_manifest.json"


def _file_fingerprint(path: Path) -> str:
    """Дешёвый отпечаток файла без чтения содержимого: размер + mtime."""
    st = path.stat()
    return f"{st.st_size}:{int(st.st_mtime)}"


def _load_manifest() -> dict:
    try:
        return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    except (FileNotFoundError, ValueError):
        return {}


def _save_manifest(manifest: dict) -> None:
    try:
        MANIFEST_PATH.write_text(
            json.dumps(manifest, ensure_ascii=False), encoding="utf-8"
        )
    except OSError as exc:
        logger.warning(f"Не удалось сохранить манифест индексации: {exc}")


# Файл-флаг статуса индексации: пишется в конце каждого прогона, чтобы факт
# завершения можно было увидеть без терминала и без активной сессии.
STATUS_PATH = config.LOG_DIR / "ingest_status.txt"


def _write_status(text: str) -> None:
    try:
        STATUS_PATH.write_text(text, encoding="utf-8")
    except OSError as exc:
        logger.warning(f"Не удалось записать статус-файл индексации: {exc}")


def list_supported_files(doc_dir: Path) -> list[Path]:
    """Все файлы поддерживаемых форматов в директории (рекурсивно, без дублей)."""
    all_files: list[Path] = []
    for ext in config.SUPPORTED_EXTENSIONS:
        all_files.extend(doc_dir.rglob(f"*{ext}"))
    return sorted(set(all_files))


def iter_file_chunks(file_path: Path) -> Generator[dict, None, None]:
    """Извлекает и нарезает ОДИН файл на чанки с метаданными."""
    pages = extract_document(file_path)
    if not pages:
        logger.warning(f"Текст не извлечён из: {file_path.name}")
        return

    file_ext = file_path.suffix.lower().lstrip(".")
    for page in pages:
        chunks = chunk_text(page["text"])
        for chunk_idx, chunk_text_val in enumerate(chunks):
            lang = detect_language(chunk_text_val)
            # Добавляем префикс для multilingual-e5
            prefixed = f"passage: {chunk_text_val}"
            yield {
                "id": make_doc_id(file_path.name, page["page_num"], chunk_idx),
                "text": prefixed,
                "metadata": {
                    "file_name": file_path.name,
                    "file_path": str(file_path),
                    "file_type": file_ext,
                    "page_num": page["page_num"],
                    "chunk_idx": chunk_idx,
                    "language": lang,
                    "char_count": len(chunk_text_val),
                },
            }


def iter_doc_chunks(doc_dir: Path) -> Generator[dict, None, None]:
    """
    Рекурсивно обходит директорию, обрабатывает все поддерживаемые форматы.
    Для каждого чанка возвращает словарь с текстом и метаданными.
    """
    all_files = list_supported_files(doc_dir)
    logger.info(f"Найдено файлов для индексации: {len(all_files)}")
    if not all_files:
        logger.warning(f"Файлы не найдены в директории: {doc_dir}")
        return

    for file_path in tqdm(all_files, desc="Обработка документов", unit="файл"):
        yield from iter_file_chunks(file_path)


# ─────────────────────────────────────────────────────────────
# Основная функция индексации
# ─────────────────────────────────────────────────────────────

def run_ingestion(doc_dir: Path, reset: bool = False) -> None:
    """Запускает полную индексацию документов из директории."""
    t0 = time.time()
    logger.info("=" * 60)
    logger.info("Запуск индексации архива документов")
    logger.info(f"Директория: {doc_dir}")
    logger.info(f"Форматы: {', '.join(sorted(config.SUPPORTED_EXTENSIONS))}")
    logger.info("=" * 60)

    # Модель эмбеддингов
    logger.info(
        f"Загрузка модели эмбеддингов: {config.EMBEDDING_MODEL} "
        f"(device={config.EMBEDDING_DEVICE or 'auto'})"
    )
    model = SentenceTransformer(config.EMBEDDING_MODEL, device=config.EMBEDDING_DEVICE)

    # ChromaDB
    collection = get_collection(reset=reset)
    existing_ids = set(collection.get(include=[])["ids"])
    logger.info(f"Уже проиндексировано чанков: {len(existing_ids)}")

    # Манифест: пропускаем парсинг неизменённых уже проиндексированных файлов.
    manifest = {} if reset else _load_manifest()
    if reset:
        _save_manifest({})

    all_files = list_supported_files(doc_dir)
    logger.info(f"Найдено файлов: {len(all_files)}")
    if not all_files:
        logger.warning(f"Файлы не найдены в директории: {doc_dir}")
        return

    # Накапливаем батч
    batch_ids, batch_texts, batch_metas = [], [], []
    new_chunks = 0
    skipped_chunks = 0
    skipped_files = 0
    parsed_fps: dict[str, str] = {}   # отпечатки обработанных в этом прогоне файлов

    for file_path in tqdm(all_files, desc="Обработка документов", unit="файл"):
        key = str(file_path)
        try:
            fp = _file_fingerprint(file_path)
        except OSError:
            fp = None

        # File-level skip: файл не менялся и уже зафиксирован в манифесте —
        # не парсим его вовсе (главная экономия времени).
        if fp is not None and manifest.get(key) == fp:
            skipped_files += 1
            continue

        for chunk in iter_file_chunks(file_path):
            if chunk["id"] in existing_ids:
                skipped_chunks += 1
                continue

            batch_ids.append(chunk["id"])
            batch_texts.append(chunk["text"])
            batch_metas.append(chunk["metadata"])

            if len(batch_ids) >= config.CHROMA_BATCH_SIZE:
                _flush_batch(model, collection, batch_ids, batch_texts, batch_metas)
                new_chunks += len(batch_ids)
                batch_ids, batch_texts, batch_metas = [], [], []

        # Файл разобран целиком — запомним отпечаток (чанки в батче/базе).
        if fp is not None:
            parsed_fps[key] = fp

    # Остаток
    if batch_ids:
        _flush_batch(model, collection, batch_ids, batch_texts, batch_metas)
        new_chunks += len(batch_ids)

    # Сохраняем манифест один раз в конце успешного прогона.
    manifest.update(parsed_fps)
    _save_manifest(manifest)

    elapsed = time.time() - t0
    total = collection.count()
    logger.info("=" * 60)
    logger.info(f"Индексация завершена за {elapsed:.1f} сек.")
    logger.info(f"Добавлено новых чанков: {new_chunks}")
    logger.info(f"Пропущено чанков (уже в базе): {skipped_chunks}")
    logger.info(f"Пропущено файлов (не менялись): {skipped_files}")
    logger.info(f"Всего в базе: {total}")
    logger.info("=" * 60)

    _write_status(
        "✅ ИНДЕКСАЦИЯ ЗАВЕРШЕНА: {ts}\n"
        "Длительность: {dur:.0f} сек\n"
        "Добавлено новых чанков: {new}\n"
        "Пропущено файлов (не менялись): {skf}\n"
        "Всего чанков в базе: {tot}\n".format(
            ts=time.strftime("%Y-%m-%d %H:%M:%S"),
            dur=elapsed, new=new_chunks, skf=skipped_files, tot=total,
        )
    )


def _flush_batch(model, collection, ids, texts, metas):
    """Генерирует эмбеддинги и записывает батч в ChromaDB."""
    embeddings = model.encode(
        texts,
        batch_size=config.EMBEDDING_BATCH_SIZE,
        normalize_embeddings=config.NORMALIZE_EMBEDDINGS,
        show_progress_bar=False,
    ).tolist()
    collection.add(
        ids=ids,
        documents=texts,
        embeddings=embeddings,
        metadatas=metas,
    )


# ─────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────
# Синхронизация удалений: убрать из базы то, чего больше нет в архиве
# ─────────────────────────────────────────────────────────────

def remove_file(file_path, prune_notes: bool = False, dry_run: bool = False) -> dict:
    """
    Точечно удаляет данные ОДНОГО файла: чанки из ChromaDB (по `file_path`),
    запись из манифеста и, опционально, его заметку Obsidian (по `source_file`).
    Используется наблюдателем при событии удаления файла.
    """
    fp = str(file_path)
    collection = get_collection(reset=False)
    n_before = collection.count()
    if not dry_run:
        collection.delete(where={"file_path": fp})
    n_deleted = n_before - collection.count() if not dry_run else 0

    if not dry_run:
        manifest = _load_manifest()
        if fp in manifest:
            del manifest[fp]
            _save_manifest(manifest)

    removed_note = False
    if prune_notes:
        name = Path(file_path).name
        vault = Path(config.KMS_VAULT_DIR)
        if vault.exists():
            for md in vault.glob("*.md"):
                head = md.read_text(encoding="utf-8", errors="ignore")
                m = re.search(r'^source_file:\s*"?(.+?)"?\s*$', head, re.MULTILINE)
                if m and m.group(1).strip() == name:
                    if not dry_run:
                        md.unlink()
                    removed_note = True
                    break

    logger.info(
        "Удаление '%s': чанков убрано %d, заметка %s%s",
        name if prune_notes else Path(fp).name, n_deleted,
        "удалена" if removed_note else "не трогалась",
        " [DRY-RUN]" if dry_run else "",
    )
    return {"chunks": n_deleted, "note_removed": removed_note}


def prune_deleted(doc_dir: Path, prune_notes: bool = False, dry_run: bool = False) -> dict:
    """
    Полная сверка с архивом: убирает из индекса (и опц. из заметок) данные файлов,
    которых больше нет на диске. Чанки — по метаданным `file_path`, записи
    манифеста — по ключу, заметки — по `source_file` во фронтматтере.

    ЗАЩИТА: если архив не существует или пуст (например, отмонтирован внешний
    диск), prune ОТМЕНЯЕТСЯ — иначе пустой список «текущих файлов» снёс бы весь
    индекс.
    """
    t0 = time.time()
    logger.info("=" * 60)
    logger.info("Синхронизация удалений (prune)%s", " [DRY-RUN]" if dry_run else "")
    logger.info(f"Директория: {doc_dir}")

    current_files = list_supported_files(doc_dir)
    current_paths = {str(p) for p in current_files}
    current_names = {p.name for p in current_files}

    if not Path(doc_dir).exists() or not current_paths:
        logger.error(
            "Архив пуст или недоступен (%s). Prune ОТМЕНЁН, чтобы не удалить весь "
            "индекс. Проверьте, примонтирован ли диск.", doc_dir,
        )
        return {"aborted": True, "orphan_chunks": 0, "orphan_files": 0, "removed_notes": 0}

    collection = get_collection(reset=False)
    data = collection.get(include=["metadatas"])
    ids = data["ids"]
    metas = data["metadatas"] or []

    orphan_ids, orphan_paths = [], set()
    for cid, meta in zip(ids, metas):
        fp = (meta or {}).get("file_path")
        if fp not in current_paths:
            orphan_ids.append(cid)
            orphan_paths.add(fp)

    logger.info(
        "Чанков всего: %d | осиротевших файлов: %d | осиротевших чанков: %d",
        len(ids), len(orphan_paths), len(orphan_ids),
    )

    if not dry_run and orphan_ids:
        for i in range(0, len(orphan_ids), 500):
            collection.delete(ids=orphan_ids[i:i + 500])

    removed_manifest = 0
    if not dry_run:
        manifest = _load_manifest()
        kept = {k: v for k, v in manifest.items() if k in current_paths}
        removed_manifest = len(manifest) - len(kept)
        _save_manifest(kept)

    removed_notes = 0
    if prune_notes:
        vault = Path(config.KMS_VAULT_DIR)
        if vault.exists():
            for md in vault.glob("*.md"):
                head = md.read_text(encoding="utf-8", errors="ignore")
                m = re.search(r'^source_file:\s*"?(.+?)"?\s*$', head, re.MULTILINE)
                if m and m.group(1).strip() not in current_names:
                    if not dry_run:
                        md.unlink()
                    removed_notes += 1
        logger.info("Осиротевших заметок: %d", removed_notes)

    elapsed = time.time() - t0
    logger.info(
        "Prune завершён за %.1f сек%s. Чанков: %d, файлов: %d, заметок: %d, "
        "записей манифеста: %d.",
        elapsed, " [DRY-RUN]" if dry_run else "",
        0 if dry_run else len(orphan_ids), len(orphan_paths),
        0 if dry_run else removed_notes, removed_manifest,
    )
    logger.info("=" * 60)
    return {
        "aborted": False,
        "orphan_chunks": len(orphan_ids),
        "orphan_files": len(orphan_paths),
        "removed_notes": removed_notes,
        "dry_run": dry_run,
    }


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Индексация документов (PDF, DOCX, TXT и др.) в ChromaDB"
    )
    parser.add_argument(
        "--doc-dir", "--pdf-dir",   # --pdf-dir оставлен для обратной совместимости
        type=Path,
        default=config.KMS_ARCHIVE_DIR,
        help=f"Путь к директории с документами (по умолчанию: {config.KMS_ARCHIVE_DIR})",
    )
    parser.add_argument(
        "--reset",
        action="store_true",
        help="Полностью сбросить базу ChromaDB перед индексацией",
    )
    parser.add_argument(
        "--prune",
        action="store_true",
        help="После индексации убрать из базы данные файлов, удалённых из архива",
    )
    parser.add_argument(
        "--prune-only",
        action="store_true",
        help="Только синхронизация удалений (без индексации)",
    )
    parser.add_argument(
        "--prune-notes",
        action="store_true",
        help="При prune также удалять осиротевшие заметки Obsidian",
    )
    parser.add_argument(
        "--remove-file",
        type=Path,
        default=None,
        help="Точечно удалить данные одного файла (путь); используется наблюдателем",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Только показать, что будет удалено при prune/remove-file, ничего не удаляя",
    )
    args = parser.parse_args()

    if args.remove_file is not None:
        remove_file(args.remove_file, prune_notes=args.prune_notes, dry_run=args.dry_run)
    elif args.prune_only:
        prune_deleted(doc_dir=args.doc_dir, prune_notes=args.prune_notes, dry_run=args.dry_run)
    else:
        run_ingestion(doc_dir=args.doc_dir, reset=args.reset)
        if args.prune:
            prune_deleted(doc_dir=args.doc_dir, prune_notes=args.prune_notes, dry_run=args.dry_run)
